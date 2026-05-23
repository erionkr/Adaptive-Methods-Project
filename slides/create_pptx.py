"""
Create PowerPoint presentation for CO Project – Adaptive Algorithms
Design: Uni Basel style (teal header bar, serif font, minimal text, figure-heavy)
Matching Chapter 22 Deep Learning presentation design.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colours ──────────────────────────────────────────────────────
TEAL = RGBColor(160, 196, 184)      # Uni Basel teal
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(80, 80, 80)
TEAL_DARK = RGBColor(100, 150, 135)

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)
FIGURES = os.path.join(os.path.dirname(__file__), '..', 'figures')

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper functions ─────────────────────────────────────────────

def add_teal_header(slide, title_text, full_bg=False):
    """Add teal header bar at top (or full background for section slides)."""
    if full_bg:
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = TEAL
    else:
        # Header bar
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0), Inches(0),
                                      SLIDE_W, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()
        # Title text inside bar
        tf = bar.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title_text
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK
        tf.paragraphs[0].font.name = "Georgia"
        tf.margin_left = Inches(0.5)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_footer(slide, authors="Erion Krasniqi and Denis Mustafa Xhabrahimi"):
    """Add footer line + text."""
    # Black line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.4), Inches(7.05),
                                   Inches(12.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()
    # Author text left
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(6), Inches(0.3))
    tf = txBox.text_frame
    tf.paragraphs[0].text = authors
    tf.paragraphs[0].font.size = Pt(9)
    tf.paragraphs[0].font.color.rgb = DARK_GRAY
    tf.paragraphs[0].font.name = "Georgia"
    # Uni Basel right
    txBox2 = slide.shapes.add_textbox(Inches(10), Inches(7.1), Inches(3), Inches(0.3))
    tf2 = txBox2.text_frame
    tf2.paragraphs[0].text = "Universität Basel"
    tf2.paragraphs[0].font.size = Pt(9)
    tf2.paragraphs[0].font.color.rgb = DARK_GRAY
    tf2.paragraphs[0].font.name = "Georgia"
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_bullet_slide(title, bullets, speaker=""):
    """Content slide with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_teal_header(slide, title)
    add_footer(slide)
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.font.size = Pt(22)
        p.font.name = "Georgia"
        p.font.color.rgb = BLACK
        p.space_after = Pt(14)
        p.level = 0
        # Sub-bullets: start with "  •"
        if b.startswith("  "):
            p.level = 1
            p.font.size = Pt(18)
    # Speaker note
    if speaker:
        txBox3 = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.3))
        tf3 = txBox3.text_frame
        tf3.paragraphs[0].text = f"🎤 {speaker}"
        tf3.paragraphs[0].font.size = Pt(10)
        tf3.paragraphs[0].font.color.rgb = DARK_GRAY
        tf3.paragraphs[0].font.name = "Georgia"
        tf3.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide

def add_section_slide(title):
    """Teal full-background section divider."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, "", full_bg=True)
    add_footer(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
    tf = txBox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = BLACK
    tf.paragraphs[0].font.name = "Georgia"
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

def add_figure_slide(title, fig_paths, caption="", speaker=""):
    """Slide with 1 or 2 figures side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title)
    add_footer(slide)
    if len(fig_paths) == 1:
        fpath = os.path.join(FIGURES, fig_paths[0])
        if os.path.exists(fpath):
            slide.shapes.add_picture(fpath, Inches(2.5), Inches(1.3), Inches(8), Inches(4.8))
    elif len(fig_paths) == 2:
        for i, fp in enumerate(fig_paths):
            fpath = os.path.join(FIGURES, fp)
            if os.path.exists(fpath):
                left = Inches(0.5) if i == 0 else Inches(6.8)
                slide.shapes.add_picture(fpath, left, Inches(1.3), Inches(6), Inches(3.8))
    elif len(fig_paths) == 3:
        for i, fp in enumerate(fig_paths):
            fpath = os.path.join(FIGURES, fp)
            if os.path.exists(fpath):
                left = Inches(0.3 + i * 4.2)
                slide.shapes.add_picture(fpath, left, Inches(1.3), Inches(4), Inches(3.5))
    # Caption
    if caption:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = caption
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.color.rgb = DARK_GRAY
        tf.paragraphs[0].font.name = "Georgia"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if speaker:
        txBox3 = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.3))
        tf3 = txBox3.text_frame
        tf3.paragraphs[0].text = f"🎤 {speaker}"
        tf3.paragraphs[0].font.size = Pt(10)
        tf3.paragraphs[0].font.color.rgb = DARK_GRAY
        tf3.paragraphs[0].font.name = "Georgia"
        tf3.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide

def add_math_slide(title, lines, speaker=""):
    """Slide with centered math/algorithm text (for update rules)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title)
    add_footer(slide)
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.name = "Georgia"
        p.font.color.rgb = BLACK
        p.space_after = Pt(10)
    if speaker:
        txBox3 = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.3))
        tf3 = txBox3.text_frame
        tf3.paragraphs[0].text = f"🎤 {speaker}"
        tf3.paragraphs[0].font.size = Pt(10)
        tf3.paragraphs[0].font.color.rgb = DARK_GRAY
        tf3.paragraphs[0].font.name = "Georgia"
        tf3.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide

def add_table_slide(title, headers, rows, speaker=""):
    """Slide with a table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title)
    add_footer(slide)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table = slide.shapes.add_table(n_rows, n_cols,
                                    Inches(0.8), Inches(1.5),
                                    Inches(11.5), Inches(0.45 * n_rows)).table
    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.name = "Georgia"
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Georgia"
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(240, 248, 245)
    if speaker:
        txBox3 = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.3))
        tf3 = txBox3.text_frame
        tf3.paragraphs[0].text = f"🎤 {speaker}"
        tf3.paragraphs[0].font.size = Pt(10)
        tf3.paragraphs[0].font.color.rgb = DARK_GRAY
        tf3.paragraphs[0].font.name = "Georgia"
        tf3.paragraphs[0].alignment = PP_ALIGN.RIGHT
    return slide


# =====================================================================
# SLIDES  —  ~20 main slides for 15-20 min + backup slides
# =====================================================================

# ── 1. Title Slide ───────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_teal_header(slide, "", full_bg=True)
txLogo = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(3), Inches(0.6))
tf = txLogo.text_frame
p = tf.paragraphs[0]
p.text = "Universität\nBasel"
p.font.size = Pt(14)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = BLACK
txTitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1.5))
tf = txTitle.text_frame
p = tf.paragraphs[0]
p.text = "Adaptive Algorithms"
p.font.size = Pt(48)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = BLACK
p2 = tf.add_paragraph()
p2.text = "Continuous Optimization — Project Presentation"
p2.font.size = Pt(24)
p2.font.name = "Georgia"
p2.font.color.rgb = DARK_GRAY
txAuth = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(1))
tf = txAuth.text_frame
p = tf.paragraphs[0]
p.text = "Erion Krasniqi and Denis Mustafa Xhabrahimi"
p.font.size = Pt(18)
p.font.italic = True
p.font.name = "Georgia"
p2 = tf.add_paragraph()
p2.text = "Frühjahrssemester / Spring Semester 2026"
p2.font.size = Pt(14)
p2.font.name = "Georgia"
p2.font.color.rgb = DARK_GRAY
add_footer(slide)

# ── 2. Outline ───────────────────────────────────────────────────
add_bullet_slide("Outline", [
    "1    Introduction & Motivation",
    "2    Related Work: SGD & Adaptive Methods",
    "3    Methodology & Results",
    "4    Discussion & Next Steps",
])

# ══════════════════════════════════════════════════════════════════
# SECTION 1: Introduction  (~2 slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Introduction")

add_bullet_slide("Why Adaptive Step Sizes?", [
    "• Gradient descent:  w_{t+1} = w_t − η ∇f(w_t)",
    "• Fixed step size η — one size does NOT fit all",
    "• Too large → divergence; too small → slow convergence",
    "",
    "• Key idea: let the algorithm choose η_t",
    "  based on the observed gradient history",
], speaker="Denis")

add_bullet_slide("Project Goal", [
    "• Survey adaptive first-order methods:",
    "  Adagrad, AdaGrad-Norm, RMSprop, Adam",
    "",
    "• Compare against vanilla SGD (baseline)",
    "",
    "• Two datasets: A9a (binary) and MNIST (multi-class)",
    "• Two losses: convex logistic + non-convex regularizer",
], speaker="Denis")

# ══════════════════════════════════════════════════════════════════
# SECTION 2: Related Work  (~6 slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Related Work")

add_math_slide("SGD — Stochastic Gradient Descent", [
    "Update rule:",
    "    w_{t+1} = w_t − η g_t ,    g_t = ∇f_{i_t}(w_t)",
    "",
    "  ✓  Simple, well-understood: O(1/√T) convergence",
    "  ✗  Requires careful tuning of η",
    "  ✗  No per-coordinate adaptation",
    "",
    "→ Baseline for all comparisons",
], speaker="Denis")

add_math_slide("Adagrad  (Duchi, Hazan & Singer, 2011)", [
    "Accumulate squared gradients per coordinate:",
    "    G_t = G_{t−1} + g_t ⊙ g_t",
    "",
    "Update:   w_{t+1} = w_t − η / (√G_t + ε) · g_t",
    "",
    "  ✓  Per-coordinate scaling — great for sparse features",
    "  ✗  Step size monotonically decreases → can stop learning",
], speaker="Denis")

add_math_slide("AdaGrad-Norm · RMSprop", [
    "AdaGrad-Norm (Ward, Wu & Bottou, 2020):",
    "    b_t² = b_{t−1}² + ‖g_t‖²     →  w_{t+1} = w_t − η/(b_t+ε) · g_t",
    "  ✓  Sharp O(1/√T) even non-convex  ✗  No per-coord adapt.",
    "",
    "RMSprop (Tieleman & Hinton, 2012):",
    "    v_t = ρ v_{t−1} + (1−ρ) g_t²  →  w_{t+1} = w_t − η/(√v_t+ε) · g_t",
    "  ✓  Forgets old gradients → adapts faster than Adagrad",
    "  ✗  No bias correction, no formal convergence guarantees",
], speaker="Erion")

add_math_slide("Adam  (Kingma & Ba, 2014)", [
    "First moment:   m_t = β₁ m_{t−1} + (1−β₁) g_t",
    "Second moment:  v_t = β₂ v_{t−1} + (1−β₂) g_t²",
    "",
    "Bias correction:",
    "    m̂_t = m_t / (1 − β₁ᵗ)  ,   v̂_t = v_t / (1 − β₂ᵗ)",
    "",
    "Update:  w_{t+1} = w_t − η m̂_t / (√v̂_t + ε)",
    "",
    "  ✓  Combines momentum + adaptive step size",
    "  ✗  May not converge to optimum (Reddi et al. 2018)",
], speaker="Denis")

add_table_slide("Comparison at a Glance",
    ["Optimizer", "Per-coord", "Momentum", "Forgetting", "Convergence"],
    [
        ["SGD",          "✗", "✗", "—", "O(1/√T)"],
        ["Adagrad",      "✓", "✗", "✗", "O(log T/√T)"],
        ["AdaGrad-Norm", "✗", "✗", "✗", "O(1/√T)"],
        ["RMSprop",      "✓", "✗", "✓", "heuristic"],
        ["Adam",         "✓", "✓", "✓", "O(1/√T)*"],
    ],
    speaker="Both"
)

# ══════════════════════════════════════════════════════════════════
# SECTION 3: Methodology & Results  (~7 slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Methodology & Results")

add_bullet_slide("Experimental Setup", [
    "Datasets:",
    "  • A9a: binary, 32K samples, 123 features",
    "  • MNIST: 10 classes, 60K samples, 780 features",
    "",
    "Protocol:",
    "  • All optimizers implemented from scratch in NumPy",
    "  • Mini-batch SGD (batch 128/256), 10–15 epochs",
    "  • Averaged over 3 seeds  (shaded = ±1 std)",
    "",
    "Losses: convex logistic  +  non-convex regularizer  r(w) = λ Σ α wⱼ²/(1+α wⱼ²)",
], speaker="Erion")

add_figure_slide("A9a — Logistic Loss (Convex)",
    ["fig01_a9a__logistic_loss_(convex).png", "fig02_a9a__test_accuracy.png"],
    "RMSprop & Adagrad converge fastest; all reach ≈85% test accuracy.",
    speaker="Erion")

add_figure_slide("A9a — Non-convex Regularizer",
    ["fig03_a9a__logistic__non-convex_regularizer.png", "fig04_a9a__test_accuracy_(non-convex).png"],
    "Non-convex regularizer raises final loss but optimizer ranking stays the same.",
    speaker="Erion")

add_figure_slide("MNIST — Softmax Cross-Entropy",
    ["fig05_mnist__softmax_cross-entropy.png", "fig06_mnist__test_accuracy.png"],
    "Adagrad dominates; RMSprop & Adam close behind. SGD needs 3× more iterations.",
    speaker="Erion")

add_figure_slide("Gradient Norms",
    ["fig07_a9a__gradient_norm_(convex).png",
     "fig08_a9a__gradient_norm_(non-convex).png",
     "fig09_mnist__gradient_norm.png"],
    "Gradient norms reveal WHY adaptive methods help: they normalize large initial gradients.",
    speaker="Erion")

add_table_slide("Summary of Results",
    ["Optimizer", "A9a Loss", "A9a Acc", "A9a NC Loss", "A9a NC Acc", "MNIST Loss", "MNIST Acc"],
    [
        ["SGD",          "0.337", "0.847", "0.349", "0.847", "0.286", "0.923"],
        ["Adagrad",      "0.328", "0.850", "0.340", "0.850", "0.271", "0.925"],
        ["AdaGrad-Norm", "0.333", "0.848", "0.344", "0.848", "0.302", "0.917"],
        ["RMSprop",      "0.328", "0.849", "0.340", "0.849", "0.278", "0.924"],
        ["Adam",         "0.330", "0.849", "0.341", "0.849", "0.276", "0.924"],
    ],
    speaker="Erion"
)

# ══════════════════════════════════════════════════════════════════
# SECTION 4: Discussion & Next Steps  (~2 slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Discussion & Next Steps")

add_bullet_slide("Key Takeaways", [
    "• Adaptive methods converge significantly faster in early iterations",
    "",
    "• RMSprop is surprisingly effective — fast forgetting helps",
    "",
    "• Final accuracy differences are small on these tasks",
    "  → adaptive methods shine in speed, not final quality",
    "",
    "• Non-convex regularizer has limited impact at moderate λ",
], speaker="Both")

add_bullet_slide("Next Steps (for the Report)", [
    "• Hyperparameter sensitivity analysis (learning rate sweeps)",
    "• Theoretical convergence rate comparison",
    "• Wall-clock time comparison (not just iterations)",
    "• Explore larger λ for non-convex regularizer",
    "• Final report in NeurIPS format (deadline: 21.07.2026)",
], speaker="Both")

# ── References ───────────────────────────────────────────────────
add_bullet_slide("References", [
    "• Kingma & Ba (2014). Adam: A method for stochastic optimization.",
    "• Duchi, Hazan & Singer (2011). Adaptive subgradient methods. JMLR 12.",
    "• Ward, Wu & Bottou (2020). AdaGrad stepsizes. JMLR 21.",
    "• Tieleman & Hinton (2012). Lecture 6.5 — RMSprop. Coursera.",
    "• Ruder (2016). An overview of gradient descent optimization.",
    "• Reddi et al. (2018). On the convergence of Adam and beyond. ICLR.",
])

# ── Thank You ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_teal_header(slide, "", full_bg=True)
add_footer(slide)
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "Thank you!"
p.font.size = Pt(48)
p.font.bold = True
p.font.name = "Georgia"
p.font.color.rgb = BLACK
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Questions?"
p2.font.size = Pt(30)
p2.font.name = "Georgia"
p2.font.color.rgb = DARK_GRAY
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(30)

# ══════════════════════════════════════════════════════════════════
# BACKUP SLIDES
# ══════════════════════════════════════════════════════════════════
add_section_slide("Backup Slides")

add_figure_slide("Learning Rate Sensitivity",
    ["fig16_lr_sweep_a9a.png"],
    "Adagrad & AdaGrad-Norm are most robust to η; RMSprop diverges at large η.",
    speaker="Erion")

add_figure_slide("Wall-Clock Time Comparison",
    ["fig11_wallclock_a9a_—_loss_vs_wall-clock_time.png",
     "fig12_wallclock_mnist_—_loss_vs_wall-clock_tim.png"],
    "Adaptive methods have slightly more per-step overhead but reach low loss faster.",
    speaker="Erion")

add_figure_slide("Non-convex Regularizer — λ Sweep",
    ["fig17_lambda_sweep.png"],
    "At λ ≥ 0.1 the regularizer dominates, pushing weights to zero.",
    speaker="Erion")

add_figure_slide("Empirical vs. Theoretical Convergence",
    ["fig14_convergence_rate_a9a_—_suboptimality_vs_it.png",
     "fig15_convergence_rate_mnist_—_suboptimality_vs_.png"],
    "All methods converge faster than worst-case O(1/√T). Adagrad has steepest empirical rate.",
    speaker="Erion")

# ── Save ─────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), 'Adaptive_Algorithms_CO_Presentation_v2.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
