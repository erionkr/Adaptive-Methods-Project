"""
Create PowerPoint presentation for CO Project – Adaptive Algorithms
Design: Uni Basel style (teal header bar, Georgia font, figure-heavy)
Math formulas are rendered as PNG images via matplotlib for proper typesetting.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from io import BytesIO
import os, hashlib

# ── Colours ──────────────────────────────────────────────────────
TEAL  = RGBColor(160, 196, 184)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY  = RGBColor(80, 80, 80)
GREEN_DARK = RGBColor(40, 120, 40)
RED_DARK   = RGBColor(180, 50, 50)

SLIDE_W = Inches(13.333)   # 16:9
SLIDE_H = Inches(7.5)
FIGURES = os.path.join(os.path.dirname(__file__), '..', 'figures')
FORMULA_DIR = os.path.join(os.path.dirname(__file__), '_formulas')
os.makedirs(FORMULA_DIR, exist_ok=True)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ── Formula rendering ────────────────────────────────────────────

def render_formula(latex, fontsize=22, dpi=200):
    """Render a LaTeX formula to a PNG file, return its path.
    Uses matplotlib mathtext (no LaTeX install needed)."""
    # Cache by hash so we don't re-render
    key = hashlib.md5(f"{latex}_{fontsize}".encode()).hexdigest()[:12]
    path = os.path.join(FORMULA_DIR, f"f_{key}.png")
    if os.path.exists(path):
        return path
    fig, ax = plt.subplots(figsize=(0.01, 0.01))
    ax.axis('off')
    fig.patch.set_alpha(0.0)
    text = ax.text(0, 0.5, f"${latex}$",
                   fontsize=fontsize,
                   transform=ax.transAxes,
                   verticalalignment='center',
                   color='black',
                   fontfamily='serif')
    fig.savefig(path, dpi=dpi, bbox_inches='tight',
                pad_inches=0.05, transparent=True)
    plt.close(fig)
    return path

# ── Layout helpers ───────────────────────────────────────────────

def add_teal_header(slide, title_text, full_bg=False):
    if full_bg:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = TEAL
    else:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0), Inches(0),
                                      SLIDE_W, Inches(1.1))
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL
        bar.line.fill.background()
        tf = bar.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = title_text
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BLACK
        tf.paragraphs[0].font.name = "Georgia"
        tf.margin_left = Inches(0.5)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_footer(slide, authors="Erion Krasniqi and Denis Mustafa Xhabrahimi"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.4), Inches(7.05),
                                   Inches(12.5), Pt(1))
    line.fill.solid()
    line.fill.fore_color.rgb = BLACK
    line.line.fill.background()
    tb1 = slide.shapes.add_textbox(Inches(0.4), Inches(7.1), Inches(6), Inches(0.3))
    p = tb1.text_frame.paragraphs[0]
    p.text = authors; p.font.size = Pt(9); p.font.color.rgb = DARK_GRAY; p.font.name = "Georgia"
    tb2 = slide.shapes.add_textbox(Inches(10), Inches(7.1), Inches(3), Inches(0.3))
    p = tb2.text_frame.paragraphs[0]
    p.text = "Universität Basel"; p.font.size = Pt(9); p.font.color.rgb = DARK_GRAY
    p.font.name = "Georgia"; p.alignment = PP_ALIGN.RIGHT

def _add_speaker(slide, speaker):
    if not speaker:
        return
    tb = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = f"🎤 {speaker}"; p.font.size = Pt(10)
    p.font.color.rgb = DARK_GRAY; p.font.name = "Georgia"; p.alignment = PP_ALIGN.RIGHT

def _text_line(tf, text, size=22, color=BLACK, bold=False, space_after=12, italic=False):
    """Append a text paragraph to an existing text frame."""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text; p.font.size = Pt(size); p.font.name = "Georgia"
    p.font.color.rgb = color; p.font.bold = bold; p.font.italic = italic
    p.space_after = Pt(space_after)
    return p

# ── Slide builders ───────────────────────────────────────────────

def add_section_slide(title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, "", full_bg=True)
    add_footer(slide)
    tb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
    p = tb.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(44); p.font.bold = True
    p.font.color.rgb = BLACK; p.font.name = "Georgia"; p.alignment = PP_ALIGN.CENTER

def add_bullet_slide(title, bullets, speaker=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title); add_footer(slide); _add_speaker(slide, speaker)
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.2))
    tf = tb.text_frame; tf.word_wrap = True
    for b in bullets:
        _text_line(tf, b, size=20 if b.startswith("  ") else 22,
                   color=DARK_GRAY if b.startswith("  ") else BLACK)
    return slide

def add_figure_slide(title, fig_paths, caption="", speaker=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title); add_footer(slide); _add_speaker(slide, speaker)
    n = len(fig_paths)
    for i, fp in enumerate(fig_paths):
        fpath = os.path.join(FIGURES, fp)
        if not os.path.exists(fpath):
            continue
        if n == 1:
            slide.shapes.add_picture(fpath, Inches(2.5), Inches(1.3), Inches(8), Inches(4.8))
        elif n == 2:
            left = Inches(0.5) if i == 0 else Inches(6.8)
            slide.shapes.add_picture(fpath, left, Inches(1.3), Inches(6), Inches(3.8))
        elif n == 3:
            slide.shapes.add_picture(fpath, Inches(0.3 + i*4.2), Inches(1.3), Inches(4), Inches(3.5))
    if caption:
        tb = slide.shapes.add_textbox(Inches(1), Inches(6.3), Inches(11), Inches(0.6))
        tf = tb.text_frame; tf.word_wrap = True
        _text_line(tf, caption, size=14, color=DARK_GRAY)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return slide

def add_table_slide(title, headers, rows, speaker=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title); add_footer(slide); _add_speaker(slide, speaker)
    nr, nc = len(rows)+1, len(headers)
    tbl = slide.shapes.add_table(nr, nc, Inches(0.8), Inches(1.5),
                                  Inches(11.5), Inches(0.45*nr)).table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs: p.font.size=Pt(14); p.font.bold=True; p.font.name="Georgia"
        c.fill.solid(); c.fill.fore_color.rgb = TEAL
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i+1, j); c.text = str(val)
            for p in c.text_frame.paragraphs: p.font.size=Pt(13); p.font.name="Georgia"
            if i % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = RGBColor(240,248,245)
    return slide

def add_optimizer_slide(title, items, speaker=""):
    """Slide with mixed text + rendered formula images.
    `items` is a list of dicts:
      {"type": "text",    "text": "...", "color": BLACK, "bold": False}
      {"type": "formula", "latex": "...", "height": 0.5}
      {"type": "pro",     "text": "..."}   # green +
      {"type": "con",     "text": "..."}   # red  --
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_teal_header(slide, title); add_footer(slide); _add_speaker(slide, speaker)

    y = 1.35  # starting y position in inches
    left = 0.7
    max_w = 11.5

    for item in items:
        if item["type"] == "formula":
            latex = item["latex"]
            h_in  = item.get("height", 0.5)
            fpath = render_formula(latex, fontsize=item.get("fontsize", 22))
            # Get image aspect ratio to compute width
            from PIL import Image
            img = Image.open(fpath)
            aspect = img.width / img.height
            w_in = min(h_in * aspect, max_w)
            slide.shapes.add_picture(fpath, Inches(left + 0.3), Inches(y),
                                      Inches(w_in), Inches(h_in))
            y += h_in + 0.1

        elif item["type"] in ("text", "pro", "con"):
            tb = slide.shapes.add_textbox(Inches(left), Inches(y), Inches(max_w), Inches(0.45))
            tf = tb.text_frame; tf.word_wrap = True
            txt = item["text"]
            color = BLACK
            if item["type"] == "pro":
                txt = "✓  " + txt; color = GREEN_DARK
            elif item["type"] == "con":
                txt = "✗  " + txt; color = RED_DARK
            elif "color" in item:
                color = item["color"]
            p = tf.paragraphs[0]
            p.text = txt; p.font.size = Pt(item.get("size", 20)); p.font.name = "Georgia"
            p.font.color.rgb = color; p.font.bold = item.get("bold", False)
            y += 0.42

        elif item["type"] == "spacer":
            y += item.get("height", 0.2)

    return slide


# =====================================================================
# SLIDES  — 15 min main + backup for Q&A
# =====================================================================

# ── 1. TITLE SLIDE ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_teal_header(slide, "", full_bg=True)
tb = slide.shapes.add_textbox(Inches(0.4), Inches(0.3), Inches(3), Inches(0.6))
p = tb.text_frame.paragraphs[0]
p.text = "Universität\nBasel"; p.font.size = Pt(14); p.font.bold = True
p.font.name = "Georgia"; p.font.color.rgb = BLACK

tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1.5))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Adaptive Algorithms"; p.font.size = Pt(48)
p.font.bold = True; p.font.name = "Georgia"; p.font.color.rgb = BLACK
p2 = tf.add_paragraph(); p2.text = "Continuous Optimization — Project Presentation"
p2.font.size = Pt(24); p2.font.name = "Georgia"; p2.font.color.rgb = DARK_GRAY

tb = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(8), Inches(1))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Erion Krasniqi and Denis Mustafa Xhabrahimi"
p.font.size = Pt(18); p.font.italic = True; p.font.name = "Georgia"
p2 = tf.add_paragraph(); p2.text = "20 / 05 / 2026"
p2.font.size = Pt(14); p2.font.name = "Georgia"; p2.font.color.rgb = DARK_GRAY
add_footer(slide)

# ── 2. OUTLINE ──────────────────────────────────────────────────
add_bullet_slide("Outline", [
    "1    Introduction & Motivation                          ~5 min",
    "2    Related Work: SGD & Adaptive Methods       ~5 min",
    "3    Methodology & Results                                ~5 min",
])

# ══════════════════════════════════════════════════════════════════
# SECTION 1 — Introduction  (~5 min, 2 content slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Introduction")

add_optimizer_slide("Why Adaptive Step Sizes?", [
    {"type": "text",    "text": "Gradient descent:", "bold": True},
    {"type": "formula", "latex": r"w_{t+1} = w_t - \eta \, \nabla f(w_t)", "height": 0.55},
    {"type": "spacer",  "height": 0.1},
    {"type": "text",    "text": "Fixed step size η — one size does not fit all"},
    {"type": "text",    "text": "Too large → divergence;   too small → slow convergence"},
    {"type": "spacer",  "height": 0.15},
    {"type": "text",    "text": "Key idea: let the algorithm choose ηₜ based on", "bold": True},
    {"type": "text",    "text": "the observed gradient history"},
], speaker="Denis")

add_bullet_slide("Project Goal", [
    "• Survey adaptive first-order methods:",
    "  Adagrad, AdaGrad-Norm, RMSprop, Adam",
    "",
    "• Compare against vanilla SGD (baseline)",
    "",
    "• Two datasets: A9a (binary) and MNIST (multi-class)",
    "• Two losses: convex logistic + non-convex regularizer",
], speaker="Denis")

# --- Pipeline graphic below Project Goal text ---
slide = prs.slides[-1]  # last added slide
pipeline_path = os.path.join(FIGURES, 'visual_pipeline.png')
if os.path.exists(pipeline_path):
    slide.shapes.add_picture(pipeline_path, Inches(1.5), Inches(4.8), Inches(10), Inches(2.0))

# ══════════════════════════════════════════════════════════════════
# SECTION 2 — Related Work  (~5 min, 6 content slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Related Work")

# --- SGD ---
add_optimizer_slide("SGD — Stochastic Gradient Descent", [
    {"type": "text",    "text": "Update rule:", "bold": True},
    {"type": "formula", "latex": r"w_{t+1} = w_t - \eta \, g_t \;,\quad g_t = \nabla f_{i_t}(w_t)", "height": 0.55},
    {"type": "spacer"},
    {"type": "pro",     "text": "Simple, well-understood convergence: O(1/√T)"},
    {"type": "con",     "text": "Requires careful tuning of η"},
    {"type": "con",     "text": "No per-coordinate adaptation"},
    {"type": "spacer"},
    {"type": "text",    "text": "→ Baseline for all comparisons", "bold": True},
], speaker="Denis")

# --- Adagrad ---
add_optimizer_slide("Adagrad  (Duchi, Hazan & Singer, 2011)", [
    {"type": "text",    "text": "Accumulate squared gradients per coordinate:", "bold": True},
    {"type": "formula", "latex": r"G_t = G_{t-1} + g_t \odot g_t", "height": 0.5},
    {"type": "spacer",  "height": 0.05},
    {"type": "text",    "text": "Update:", "bold": True},
    {"type": "formula", "latex": r"w_{t+1} = w_t - \frac{\eta}{\sqrt{G_t} + \varepsilon} \, g_t", "height": 0.6},
    {"type": "spacer"},
    {"type": "pro",     "text": "Per-coordinate scaling — great for sparse features"},
    {"type": "con",     "text": "Step size monotonically decreases → can stop learning"},
], speaker="Denis")

# --- AdaGrad-Norm + RMSprop (combined) ---
add_optimizer_slide("AdaGrad-Norm · RMSprop", [
    {"type": "text",    "text": "AdaGrad-Norm (Ward, Wu & Bottou, 2020):", "bold": True, "size": 18},
    {"type": "formula", "latex": r"b_t^2 = b_{t-1}^2 + \|g_t\|^2 \;\;\longrightarrow\;\; w_{t+1} = w_t - \frac{\eta}{b_t + \varepsilon}\,g_t", "height": 0.55},
    {"type": "pro",     "text": "Sharp O(1/√T) convergence even on non-convex landscapes"},
    {"type": "con",     "text": "No per-coordinate adaptation (single scalar)"},
    {"type": "spacer",  "height": 0.15},
    {"type": "text",    "text": "RMSprop (Tieleman & Hinton, 2012):", "bold": True, "size": 18},
    {"type": "formula", "latex": r"v_t = \rho \, v_{t-1} + (1-\rho)\,g_t^2 \;\;\longrightarrow\;\; w_{t+1} = w_t - \frac{\eta}{\sqrt{v_t}+\varepsilon}\,g_t", "height": 0.55},
    {"type": "pro",     "text": "Forgets old gradients → adapts faster than Adagrad"},
    {"type": "con",     "text": "No bias correction, no formal convergence guarantees"},
], speaker="Erion")

# --- Adam ---
add_optimizer_slide("Adam  (Kingma & Ba, 2014)", [
    {"type": "text",    "text": "First & second moment:", "bold": True},
    {"type": "formula", "latex": r"m_t = \beta_1 m_{t-1} + (1-\beta_1) g_t \;,\quad v_t = \beta_2 v_{t-1} + (1-\beta_2) g_t^2", "height": 0.55},
    {"type": "text",    "text": "Bias correction:", "bold": True},
    {"type": "formula", "latex": r"\hat{m}_t = \frac{m_t}{1 - \beta_1^t} \;,\quad \hat{v}_t = \frac{v_t}{1 - \beta_2^t}", "height": 0.6},
    {"type": "text",    "text": "Update:", "bold": True},
    {"type": "formula", "latex": r"w_{t+1} = w_t - \frac{\eta \, \hat{m}_t}{\sqrt{\hat{v}_t} + \varepsilon}", "height": 0.65},
    {"type": "pro",     "text": "Combines momentum + adaptive step size"},
    {"type": "con",     "text": "May not converge to optimal solution (Reddi et al. 2018)"},
], speaker="Denis")

# --- Optimizer Family Tree (before comparison table) ---
add_figure_slide("How the Methods Relate",
    ["visual_optimizer_tree.png"],
    caption="",
    speaker="Erion")

# --- Comparison table ---
add_table_slide("Comparison at a Glance",
    ["Optimizer", "Per-coord", "Momentum", "Forgetting", "Convergence"],
    [
        ["SGD",          "✗", "✗", "—",  "O(1/√T)"],
        ["Adagrad",      "✓", "✗", "✗",  "O(log T/√T)"],
        ["AdaGrad-Norm", "✗", "✗", "✗",  "O(1/√T)"],
        ["RMSprop",      "✓", "✗", "✓",  "heuristic"],
        ["Adam",         "✓", "✓", "✓",  "O(1/√T)*"],
    ],
    speaker="Both"
)

# ══════════════════════════════════════════════════════════════════
# SECTION 3 — Methodology & Results  (~5 min, 4 content slides)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Methodology & Results")

# --- Experimental Setup (with loss formula) ---
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_teal_header(slide, "Experimental Setup"); add_footer(slide); _add_speaker(slide, "Erion")

# Left column
tb_left = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5.8), Inches(5.5))
tf = tb_left.text_frame; tf.word_wrap = True
_text_line(tf, "Datasets", size=20, bold=True)
_text_line(tf, "• A9a: binary, 32K samples, 123 features", size=17, color=DARK_GRAY)
_text_line(tf, "• MNIST: 10 classes, 60K samples, 780 features", size=17, color=DARK_GRAY)
_text_line(tf, "", size=10)
_text_line(tf, "Loss functions", size=20, bold=True)
_text_line(tf, "• Logistic regression (convex)", size=17, color=DARK_GRAY)
_text_line(tf, "• Logistic + non-convex regularizer:", size=17, color=DARK_GRAY)

# Non-convex regularizer formula
fpath_reg = render_formula(r"r(w) = \lambda \sum_j \frac{\alpha \, w_j^2}{1 + \alpha \, w_j^2}", fontsize=20)
from PIL import Image
img = Image.open(fpath_reg)
aspect = img.width / img.height
slide.shapes.add_picture(fpath_reg, Inches(1.2), Inches(4.6), Inches(0.5 * aspect), Inches(0.5))

# Right column
tb_right = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5))
tf = tb_right.text_frame; tf.word_wrap = True
_text_line(tf, "Protocol", size=20, bold=True)
_text_line(tf, "• Implemented from scratch (NumPy)", size=17, color=DARK_GRAY)
_text_line(tf, "• Mini-batch, batch size 128/256", size=17, color=DARK_GRAY)
_text_line(tf, "• 10–15 epochs, averaged over 3 seeds", size=17, color=DARK_GRAY)
_text_line(tf, "• Shaded region = ±1 std", size=17, color=DARK_GRAY)
_text_line(tf, "", size=10)
_text_line(tf, "Optimizers", size=20, bold=True)
_text_line(tf, "• SGD, Adagrad, AdaGrad-Norm", size=17, color=DARK_GRAY)
_text_line(tf, "• RMSprop, Adam", size=17, color=DARK_GRAY)

# MNIST sample digits (small, bottom-right)
mnist_path = os.path.join(FIGURES, 'visual_mnist_samples.png')
if os.path.exists(mnist_path):
    slide.shapes.add_picture(mnist_path, Inches(7.0), Inches(5.3), Inches(5.5), Inches(0.9))

# --- A9a results ---
add_figure_slide("A9a — Logistic Loss (Convex)",
    ["fig01_a9a__logistic_loss_(convex).png", "fig02_a9a__test_accuracy.png"],
    "RMSprop & Adagrad converge fastest; all reach ≈85% test accuracy.",
    speaker="Erion")

# --- MNIST results ---
add_figure_slide("MNIST — Softmax Cross-Entropy",
    ["fig05_mnist__softmax_cross-entropy.png", "fig06_mnist__test_accuracy.png"],
    "Adagrad dominates; RMSprop & Adam close behind. SGD needs 3× more iterations.",
    speaker="Erion")

# --- Non-convex reference benchmark ---
add_figure_slide("Non-convex Reference — Rosenbrock",
    ["fig18_rosenbrock_nonconvex_reference.png"],
    "Curved non-convex valley: Adam reaches the global minimum much faster, making momentum + adaptive scaling visible.",
    speaker="Erion")

# --- Key Takeaways & Outlook ---
add_bullet_slide("Key Takeaways & Outlook", [
    "• Adaptive methods converge significantly faster in early iterations",
    "",
    "• RMSprop is surprisingly effective — fast forgetting helps",
    "",
    "• Final accuracy differences are small on these tasks",
    "  → adaptive methods shine in speed, not final quality",
    "",
    "• Weak non-convex regularizer had limited impact at moderate λ",
    "",
    "• Rosenbrock reference shows Adam's advantage on non-convex objectives",
    "",
    "Remaining for the report:",
    "  LR sweeps · wall-clock comparison · larger λ analysis · NeurIPS format (21.07)",
], speaker="Both")

# ── References ───────────────────────────────────────────────────
add_bullet_slide("References", [
    "• Kingma & Ba (2014). Adam: A method for stochastic optimization.",
    "• Duchi, Hazan & Singer (2011). Adaptive subgradient methods. JMLR 12.",
    "• Ward, Wu & Bottou (2020). AdaGrad stepsizes. JMLR 21.",
    "• Tieleman & Hinton (2012). Lecture 6.5 — RMSprop. Coursera.",
    "• Reddi et al. (2018). On the convergence of Adam and beyond. ICLR.",
    "• Rosenbrock (1960). An automatic method for finding the greatest or least value of a function.",
])

# ── Thank You ────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_teal_header(slide, "", full_bg=True); add_footer(slide)
tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf = tb.text_frame
p = tf.paragraphs[0]; p.text = "Thank you!"; p.font.size = Pt(48); p.font.bold = True
p.font.name = "Georgia"; p.font.color.rgb = BLACK; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Questions?"; p2.font.size = Pt(30)
p2.font.name = "Georgia"; p2.font.color.rgb = DARK_GRAY
p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(30)

# ══════════════════════════════════════════════════════════════════
# BACKUP SLIDES (for Q&A)
# ══════════════════════════════════════════════════════════════════
add_section_slide("Backup Slides")

add_figure_slide("A9a — Non-convex Regularizer",
    ["fig03_a9a__logistic__non-convex_regularizer.png", "fig04_a9a__test_accuracy_(non-convex).png"],
    "Non-convex regularizer raises final loss but ranking stays the same.",
    speaker="Erion")

add_figure_slide("Gradient Norms",
    ["fig07_a9a__gradient_norm_(convex).png",
     "fig08_a9a__gradient_norm_(non-convex).png",
     "fig09_mnist__gradient_norm.png"],
    "Gradient norms reveal WHY adaptive methods help: they normalize large initial gradients.",
    speaker="Erion")

add_table_slide("Summary of Results",
    ["Optimizer", "A9a Loss", "A9a Acc", "NC Loss", "NC Acc", "MNIST Loss", "MNIST Acc"],
    [
        ["SGD",          "0.337", "0.847", "0.349", "0.847", "0.286", "0.923"],
        ["Adagrad",      "0.328", "0.850", "0.340", "0.850", "0.271", "0.925"],
        ["AdaGrad-Norm", "0.333", "0.848", "0.344", "0.848", "0.302", "0.917"],
        ["RMSprop",      "0.328", "0.849", "0.340", "0.849", "0.278", "0.924"],
        ["Adam",         "0.330", "0.849", "0.341", "0.849", "0.276", "0.924"],
    ], speaker="Erion")

add_figure_slide("Learning Rate Sensitivity",
    ["fig16_lr_sweep_a9a.png"],
    "Adagrad & AdaGrad-Norm are most robust to η; RMSprop diverges at large η.",
    speaker="Erion")

add_figure_slide("Wall-Clock Time Comparison",
    ["fig11_wallclock_a9a_—_loss_vs_wall-clock_time.png",
     "fig12_wallclock_mnist_—_loss_vs_wall-clock_tim.png"],
    "Adaptive methods have slightly more per-step overhead but reach low loss faster.",
    speaker="Erion")

add_figure_slide("Non-convex Regularizer — λ Sweep",
    ["fig17_lambda_sweep.png"],
    "At λ ≥ 0.1 the regularizer dominates, pushing weights to zero.",
    speaker="Erion")

add_figure_slide("Empirical vs. Theoretical Convergence",
    ["fig14_convergence_rate_a9a_—_suboptimality_vs_it.png",
     "fig15_convergence_rate_mnist_—_suboptimality_vs_.png"],
    "All methods converge faster than worst-case O(1/√T). Adagrad has steepest empirical rate.",
    speaker="Erion")

# ── Save ─────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), 'Adaptive_Algorithms_CO_v4.pptx')
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
print(f"Formula images cached in: {FORMULA_DIR}")
